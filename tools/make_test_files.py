"""@brief Generate a set of sample files covering all supported viewers.

@details Creates a test_files/ directory with fixtures: text, code, JSON,
Markdown, HTML, CSV, XLSX, DOCX, PNG, JPEG, BMP, animated GIF, SVG, ZIP,
TAR.GZ, WAV, and PDF. Also creates a nested folder to exercise lazy tree
population. Requires Pillow, openpyxl, and python-docx.
"""

import csv
import io
import json
import math
import os
import shutil
import struct
import sys
import wave
import zipfile

OUTPUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "test_files")


def make_text_files() -> None:
    """@brief Write plain-text, code, JSON, Markdown, HTML, and CSV fixtures."""
    with open(os.path.join(OUTPUT, "hello.txt"), "w", encoding="utf-8") as handle:
        handle.write("Hello, Universal Viewer!\nLine two.\nUnicode: проверка, ±, ✓\n")

    python_source = '''"""Sample module for viewer testing."""


def fib(n):
    """Return the n-th Fibonacci number."""
    if n < 2:
        return n
    return fib(n - 1) + fib(n - 2)


if __name__ == "__main__":
    for i in range(10):
        print(fib(i))
'''
    with open(os.path.join(OUTPUT, "fib.py"), "w", encoding="utf-8") as handle:
        handle.write(python_source)

    data = {
        "name": "universal-viewer",
        "version": "1.0.0",
        "features": ["text", "images", "media"],
        "nested": {"enabled": True, "attempts": 3},
    }
    with open(os.path.join(OUTPUT, "data.json"), "w", encoding="utf-8") as handle:
        json.dump(data, handle, indent=2)

    markdown = """# Sample Document

Some **bold**, *italic*, and `inline code`.

## List

- First item
- Second item
    - Nested item

## Code

```python
print("hi")
```
"""
    with open(os.path.join(OUTPUT, "README.md"), "w", encoding="utf-8") as handle:
        handle.write(markdown)

    html_doc = """<!DOCTYPE html>
<html><head><title>Sample</title></head>
<body>
<h1>HTML Sample</h1>
<p>Paragraph with <b>bold</b> and <a href="https://example.com">a link</a>.</p>
<ul><li>One</li><li>Two</li></ul>
</body></html>
"""
    with open(os.path.join(OUTPUT, "page.html"), "w", encoding="utf-8") as handle:
        handle.write(html_doc)

    with open(os.path.join(OUTPUT, "table.csv"), "w", encoding="utf-8", newline="") as handle:
        writer = csv.writer(handle, delimiter=";")
        writer.writerow(["Name", "Value", "Enabled"])
        for index in range(50):
            writer.writerow([f"row{index}", index * 1.5, index % 2 == 0])


def make_office_files() -> None:
    """@brief Write DOCX, XLSX, and PDF fixtures using third-party libraries."""
    import docx

    document = docx.Document()
    document.add_heading("Test Document", level=1)
    document.add_paragraph("Plain paragraph with ")
    run = document.add_paragraph().add_run("bold text")
    run.bold = True
    document.add_paragraph("List item one", style="List Bullet")
    document.add_paragraph("List item two", style="List Bullet")
    table = document.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "Header A"
    table.rows[0].cells[1].text = "Header B"
    table.rows[1].cells[0].text = "Cell 1"
    table.rows[1].cells[1].text = "Cell 2"
    document.save(os.path.join(OUTPUT, "document.docx"))

    import openpyxl

    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "Numbers"
    sheet.append(["ID", "Square"])
    for index in range(1, 21):
        sheet.append([index, index * index])
    second = workbook.create_sheet("Letters")
    for row in (("A", "B"), ("C", "D")):
        second.append(row)
    workbook.save(os.path.join(OUTPUT, "report.xlsx"))

    from PIL import Image

    image = Image.new("RGB", (400, 300), (40, 90, 160))
    for x in range(400):
        for y in range(0, 300, 10):
            image.putpixel((x, y), (255, 255, 255))
    image.save(os.path.join(OUTPUT, "sample.pdf"), "PDF", resolution=72.0)


def make_image_files() -> None:
    """@brief Write PNG, JPEG, BMP, animated GIF, and SVG fixtures."""
    from PIL import Image

    gradient = Image.new("RGB", (256, 128))
    for x in range(256):
        for y in range(128):
            gradient.putpixel((x, y), (x, y * 2, (x + y) % 256))
    gradient.save(os.path.join(OUTPUT, "gradient.png"))
    gradient.save(os.path.join(OUTPUT, "gradient.jpg"), quality=90)
    gradient.save(os.path.join(OUTPUT, "gradient.bmp"))

    frames = [
        Image.new("RGB", (120, 120), color) for color in ((200, 40, 40), (40, 200, 40), (40, 40, 200))
    ]
    frames[0].save(
        os.path.join(OUTPUT, "animation.gif"),
        save_all=True,
        append_images=frames[1:],
        duration=300,
        loop=0,
    )

    svg = """<svg xmlns="http://www.w3.org/2000/svg" width="200" height="200">
  <rect width="200" height="200" fill="#204080"/>
  <circle cx="100" cy="100" r="70" fill="#ffcc00"/>
  <text x="100" y="108" font-size="20" fill="#204080" text-anchor="middle">SVG</text>
</svg>
"""
    with open(os.path.join(OUTPUT, "circle.svg"), "w", encoding="utf-8") as handle:
        handle.write(svg)


def make_pptx_file() -> None:
    """@brief Write a three-slide PPTX fixture (title, bullets, image)."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Universal Viewer"
    title_slide.placeholders[1].text = "Test presentation"

    bullets = prs.slides.add_slide(prs.slide_layouts[1])
    bullets.shapes.title.text = "Bullets"
    body = bullets.placeholders[1].text_frame
    body.text = "First point"
    second = body.add_paragraph()
    second.text = "Second point"
    second.level = 1
    third = body.add_paragraph()
    third.text = "Third point with bold run"
    third.runs[0].font.bold = True

    picture_slide = prs.slides.add_slide(prs.slide_layouts[5])
    picture_slide.shapes.title.text = "Image slide"
    picture_slide.shapes.add_picture(
        os.path.join(OUTPUT, "gradient.png"), Inches(2), Inches(2), width=Inches(5)
    )

    prs.save(os.path.join(OUTPUT, "presentation.pptx"))


def make_book_files() -> None:
    """@brief Write minimal valid EPUB and FB2 fixtures for the book viewers."""
    with open(os.path.join(OUTPUT, "gradient.png"), "rb") as handle:
        png_bytes = handle.read()

    container = """<?xml version="1.0"?>
<container version="1.0" xmlns="urn:oasis:names:tc:opendocument:xmlns:container">
  <rootfiles>
    <rootfile full-path="OEBPS/content.opf" media-type="application/oebps-package+xml"/>
  </rootfiles>
</container>
"""
    opf = """<?xml version="1.0"?>
<package xmlns="http://www.idpf.org/2007/opf" version="2.0" unique-identifier="id">
  <metadata xmlns:dc="http://purl.org/dc/elements/1.1/">
    <dc:title>Test EPUB Book</dc:title>
  </metadata>
  <manifest>
    <item id="c1" href="chapter1.xhtml" media-type="application/xhtml+xml"/>
    <item id="c2" href="chapter2.xhtml" media-type="application/xhtml+xml"/>
    <item id="img1" href="images/pic.png" media-type="image/png"/>
  </manifest>
  <spine>
    <itemref idref="c1"/>
    <itemref idref="c2"/>
  </spine>
</package>
"""
    chapter1 = """<?xml version="1.0" encoding="utf-8"?>
<html xmlns="http://www.w3.org/1999/xhtml">
<head><title>Chapter One</title></head>
<body>
<h1>Chapter One</h1>
<p>Hello <b>bold</b> and <i>italic</i> text.</p>
<img src="images/pic.png" alt="picture"/>
</body>
</html>
"""
    chapter2 = """<?xml version="1.0" encoding="utf-8"?>
<html xmlns="http://www.w3.org/1999/xhtml">
<head><title>Chapter Two</title></head>
<body>
<h1>Chapter Two</h1>
<p>The end.</p>
</body>
</html>
"""
    with zipfile.ZipFile(os.path.join(OUTPUT, "sample.epub"), "w") as archive:
        archive.writestr(
            zipfile.ZipInfo("mimetype"), "application/epub+zip", compress_type=zipfile.ZIP_STORED
        )
        archive.writestr("META-INF/container.xml", container)
        archive.writestr("OEBPS/content.opf", opf)
        archive.writestr("OEBPS/chapter1.xhtml", chapter1)
        archive.writestr("OEBPS/chapter2.xhtml", chapter2)
        archive.writestr("OEBPS/images/pic.png", png_bytes)

    import base64

    fb2 = f"""<?xml version="1.0" encoding="utf-8"?>
<FictionBook xmlns="http://www.gribuser.ru/xml/fictionbook/2.0" xmlns:l="http://www.w3.org/1999/xlink">
  <description>
    <title-info>
      <author><first-name>Ivan</first-name><last-name>Testov</last-name></author>
      <book-title>Test FB2 Book</book-title>
    </title-info>
  </description>
  <body>
    <section>
      <title><p>Chapter One</p></title>
      <p>Hello <strong>bold</strong> and <em>italic</em> text.</p>
      <empty-line/>
      <image l:href="#pic.png"/>
      <poem><stanza><v>Roses are red,</v><v>Violets are blue.</v></stanza></poem>
    </section>
  </body>
  <binary id="pic.png" content-type="image/png">{base64.b64encode(png_bytes).decode('ascii')}</binary>
</FictionBook>
"""
    with open(os.path.join(OUTPUT, "sample.fb2"), "w", encoding="utf-8") as handle:
        handle.write(fb2)


def make_syntax_files() -> None:
    """@brief Write fixtures for the extended syntax-highlighting coverage."""
    fixtures = {
        "style.css": (
            "/* theme */\nbody {\n  color: #204080;\n  margin: 0 auto;\n"
            "  font-size: 14px;\n}\n@media (max-width: 600px) {\n  body { color: red; }\n}\n"
        ),
        "theme.rasi": (
            "/* rofi theme */\n* { background-color: #202020; }\n"
            "window { border: 2px; padding: 10px; }\nelement selected { text-color: white; }\n"
        ),
        "config.jsonc": (
            "// main configuration\n{\n  \"host\": \"https://example.com\",  // endpoint\n"
            "  \"retries\": 3,          /* attempts */\n  \"debug\": false\n}\n"
        ),
        "style.css.save": "/* emacs backup of style.css */\nbody { color: black; }\n",
        "fib.nim": (
            "# Fibonacci\nproc fib(n: int): int =\n  if n < 2:\n    return n\n"
            "  return fib(n - 1) + fib(n - 2)\necho fib(10)\n"
        ),
        "analysis.r": (
            "# Simple plot\nx <- 1:10\ny <- x^2\n"
            "if (TRUE) { print(mean(y)) }\nplot(x, y)\n"
        ),
        "util.m": (
            "% MATLAB helper\nfunction s = util(x)\n"
            "% comment line\n  s = sum(x);\nend\n"
        ),
        "README.org": "* Heading one\nSome **bold** text.\n** Sub heading\n- item one\n- item two\n",
        "doc.adoc": "= Document Title\n\n== Section\n\nThis is *bold* and _italic_.\n",
        "page.man": ".TH TEST 1\n.SH NAME\ntest \\- test page\n.SH SYNOPSIS\n.B test [options]\n",
        "query.sql": (
            "-- fetch users\nSELECT id, name FROM users\n"
            "WHERE age > 21 AND active = TRUE\nORDER BY name ASC;\n"
        ),
        ".editorconfig": "root = true\n\n[*]\ncharset = utf-8\nindent_style = space\n",
        ".eslintrc": "{\"semi\": false, \"quotes\": [\"error\", \"single\"]}\n",
        ".babelrc": "{\"presets\": [\"@babel/preset-env\"]}\n",
        ".dockerignore": "node_modules\n.git\n*.log\n",
        "notes.nfo": "ÚÄÄÄÄÄÄÄÄÄ¿\n³ NFO BOX ³\nÀÄÄÄÄÄÄÄÄÄÙ\n",
        "TODO.todo": "- [ ] write tests\n- [x] ship viewer\n",
        "app.pod": "=head1 NAME\n\napp - sample pod\n\n=head1 SYNOPSIS\n\napp [options]\n",
        "guide.creole": "== Title\n\n**bold** and //italic//\n",
        "page.wiki": "= Title =\n\n'''bold''' and ''italic''\n",
        "doc.textile": "h1. Title\n\n*bold* and _italic_\n",
        "kernel.log": "2026-09-03 10:00:01 INFO boot sequence complete\n" * 5,
    }
    for name, content in fixtures.items():
        with open(os.path.join(OUTPUT, name), "w", encoding="utf-8") as handle:
            handle.write(content)
    with open(os.path.join(OUTPUT, ".gitkeep"), "w", encoding="utf-8") as handle:
        handle.write("")


def make_media_and_archive_files() -> None:
    """@brief Write WAV, ZIP, and TAR.GZ fixtures."""
    sample_rate = 22050
    duration = 1.0
    with wave.open(os.path.join(OUTPUT, "tone.wav"), "w") as wav_file:
        wav_file.setnchannels(1)
        wav_file.setsampwidth(2)
        wav_file.setframerate(sample_rate)
        frames = bytearray()
        for index in range(int(sample_rate * duration)):
            value = int(12000 * math.sin(2 * math.pi * 440 * index / sample_rate))
            frames.extend(struct.pack("<h", value))
        wav_file.writeframes(bytes(frames))

    with zipfile.ZipFile(os.path.join(OUTPUT, "archive.zip"), "w") as archive:
        archive.writestr("inner/one.txt", "First file inside the archive.\n")
        archive.writestr("inner/two.txt", "Second file.\n" * 20)

    import tarfile

    with tarfile.open(os.path.join(OUTPUT, "bundle.tar.gz"), "w:gz") as archive:
        content = "TAR member content.\n".encode("utf-8")
        buffer = io.BytesIO(content)
        info = tarfile.TarInfo(name="member.txt")
        info.size = len(content)
        archive.addfile(info, buffer)


def main() -> int:
    """@brief Regenerate the whole fixture set.

    @return Process exit code.
    """
    if os.path.isdir(OUTPUT):
        shutil.rmtree(OUTPUT)
    nested = os.path.join(OUTPUT, "nested")
    os.makedirs(nested)
    make_text_files()
    make_office_files()
    make_image_files()
    make_pptx_file()
    make_book_files()
    make_syntax_files()
    make_media_and_archive_files()
    with open(os.path.join(nested, "deep.txt"), "w", encoding="utf-8") as handle:
        handle.write("File inside a nested folder.\n")
    print(f"Fixtures written to {os.path.abspath(OUTPUT)}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
