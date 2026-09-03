"""@brief Viewer for PowerPoint presentations (.pptx) via python-pptx.

@details Shows slides one by one with prev/next navigation. Extracted
content: title, text paragraphs (with bold/italic/underline runs and
indent levels), tables, and embedded raster images. Images are decoded in
process and injected into the document through QTextDocument resources -
no files are written to disk and no external application is started.
"""

import html
import os
from itertools import islice

from PyQt6.QtCore import QUrl
from PyQt6.QtGui import QImage, QPixmap, QTextDocument
from PyQt6.QtWidgets import (
    QHBoxLayout,
    QLabel,
    QTextBrowser,
    QToolButton,
    QVBoxLayout,
    QWidget,
)

from universal_viewer.viewers.base import BaseViewer

#: @brief Maximum size of an opened presentation file (256 MB).
_MAX_FILE_BYTES = 256 * 1024 * 1024

#: @brief Maximum number of slides processed from one presentation.
_MAX_SLIDES = 300

#: @brief Maximum number of content blocks rendered per slide.
_MAX_BLOCKS_PER_SLIDE = 400

#: @brief Maximum number of table rows rendered per slide.
_MAX_TABLE_ROWS = 200

#: @brief Maximum size of an embedded image kept for rendering (20 MB).
_MAX_IMAGE_BYTES = 20 * 1024 * 1024

#: @brief Maximum rendered width of an embedded image in pixels.
_MAX_IMAGE_WIDTH = 800

#: @brief Image file extensions Qt can decode from the presentation.
_IMAGE_EXTS = frozenset({"png", "jpeg", "jpg", "gif", "bmp", "tiff", "webp"})


class PptxViewer(BaseViewer):
    """@brief Slide-by-slide viewer for PowerPoint presentations."""

    SUPPORTED_EXTENSIONS = frozenset({".pptx", ".ppsx"})

    def __init__(self, parent: QWidget | None = None) -> None:
        """@brief Construct the widget with its navigation bar.

        @param parent: Optional parent widget.
        """
        super().__init__(parent)
        self._presentation = None
        self._slide_count = 0
        self._index = 0
        self._truncated = False

        self._prev_button = QToolButton(self)
        self._prev_button.setText("< Prev")
        self._next_button = QToolButton(self)
        self._next_button.setText("Next >")
        self._slide_label = QLabel("No presentation", self)

        nav = QHBoxLayout()
        nav.setContentsMargins(4, 2, 4, 2)
        nav.addWidget(self._prev_button)
        nav.addWidget(self._next_button)
        nav.addStretch(1)
        nav.addWidget(self._slide_label)
        nav.addStretch(1)

        self._browser = QTextBrowser(self)
        self._browser.setOpenLinks(False)

        layout = QVBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)
        layout.setSpacing(0)
        layout.addLayout(nav)
        layout.addWidget(self._browser, stretch=1)

        self._prev_button.clicked.connect(self._show_previous)
        self._next_button.clicked.connect(self._show_next)

    def load(self, path: str) -> None:
        """@brief Open a presentation and display the first slide.

        @details The slide count is capped at _MAX_SLIDES; the remaining
        slides are reported as truncated instead of being rendered.

        @param path: Absolute path of the PPTX file to display.
        """
        from pptx import Presentation

        self._presentation = None
        self._slide_count = 0
        self._index = 0
        self._truncated = False
        if os.path.getsize(path) > _MAX_FILE_BYTES:
            self._browser.setHtml(
                "<p style='color:#a00'>Presentation is too large to display.</p>"
            )
            return
        try:
            presentation = Presentation(path)
            total = len(presentation.slides)
            self._slide_count = min(total, _MAX_SLIDES)
            self._truncated = total > _MAX_SLIDES
            self._presentation = presentation
        except Exception as error:
            self._browser.setHtml(
                f"<p style='color:#a00'>Could not open presentation: "
                f"{html.escape(str(error))}</p>"
            )
            return
        self._show_slide(0)

    def _show_slide(self, index: int) -> None:
        """@brief Render and display the slide with the given index.

        @param index: Zero-based slide index, clamped to the valid range.
        """
        if self._presentation is None or self._slide_count == 0:
            self._browser.setHtml("<p>(empty presentation)</p>")
            self._slide_label.setText("0 / 0")
            self._update_buttons()
            return
        index = max(0, min(index, self._slide_count - 1))
        self._index = index
        body, resources = self._render_slide(index)
        self._browser.clear()
        document = self._browser.document()
        for url, pixmap in resources:
            document.addResource(
                QTextDocument.ResourceType.ImageResource, QUrl(url), pixmap
            )
        self._browser.setHtml(body)
        suffix = f" (truncated at {_MAX_SLIDES} slides)" if self._truncated else ""
        self._slide_label.setText(f"Slide {index + 1} / {self._slide_count}{suffix}")
        self._update_buttons()

    def _render_slide(self, index: int) -> tuple[str, list[tuple[str, QPixmap]]]:
        """@brief Convert one slide into an HTML fragment with image resources.

        @param index: Zero-based slide index.
        @return Tuple (HTML body, list of (resource URL, pixmap) pairs).
        """
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        slide = self._presentation.slides[index]
        resources: list[tuple[str, QPixmap]] = []
        parts: list[str] = []
        state = {"blocks": 0, "images": 0}

        try:
            title_text = slide.shapes.title.text if slide.shapes.title is not None else ""
        except Exception:
            title_text = ""
        if title_text:
            parts.append(
                f"<h2 style='text-align:center;'>{html.escape(title_text)}</h2><hr>"
            )

        def walk(shapes) -> None:
            """@brief Recursively collect content from a shape collection.

            @param shapes: Iterable of python-pptx shape objects.
            """
            for shape in shapes:
                if state["blocks"] >= _MAX_BLOCKS_PER_SLIDE:
                    return
                try:
                    shape_type = shape.shape_type
                except Exception:
                    shape_type = None
                if shape_type == MSO_SHAPE_TYPE.GROUP:
                    walk(shape.shapes)
                    continue
                if getattr(shape, "has_text_frame", False):
                    for paragraph in shape.text_frame.paragraphs:
                        state["blocks"] += 1
                        if state["blocks"] > _MAX_BLOCKS_PER_SLIDE:
                            return
                        text = _runs_to_html(paragraph)
                        if text:
                            indent = min(int(paragraph.level), 4) * 24
                            parts.append(
                                f"<p style='margin-left:{indent}px;'>&bull; {text}</p>"
                            )
                if getattr(shape, "has_table", False) and shape.has_table:
                    state["blocks"] += 1
                    parts.append(_table_to_html(shape.table))
                image_data = self._extract_image(shape)
                if image_data is not None:
                    state["blocks"] += 1
                    url = f"pptximg://{index}/{state['images']}"
                    state["images"] += 1
                    resources.append((url, image_data[1]))
                    width, height = self._scaled_size(image_data[0])
                    parts.append(
                        f"<div style='text-align:center;'>"
                        f"<img src='{url}' width='{width}' height='{height}'></div>"
                    )

        walk(slide.shapes)
        if state["blocks"] > _MAX_BLOCKS_PER_SLIDE:
            parts.append(
                f"<p><b>[slide truncated after {_MAX_BLOCKS_PER_SLIDE} blocks]</b></p>"
            )
        return "".join(parts) or "<p>(empty slide)</p>", resources

    def _extract_image(self, shape) -> tuple[QImage, QPixmap] | None:
        """@brief Decode an embedded picture shape into a pixmap.

        @param shape: python-pptx shape that may expose an .image attribute.
        @return Tuple (decoded QImage, QPixmap) or None when not renderable.
        """
        image = getattr(shape, "image", None)
        if image is None:
            return None
        try:
            blob = image.blob
            ext = (image.ext or "").lower()
        except Exception:
            return None
        if not blob or len(blob) > _MAX_IMAGE_BYTES or ext not in _IMAGE_EXTS:
            return None
        qimage = QImage.fromData(blob)
        if qimage.isNull():
            return None
        return qimage, QPixmap.fromImage(qimage)

    @staticmethod
    def _scaled_size(pixmap: QPixmap) -> tuple[int, int]:
        """@brief Compute display size for an image, capped to _MAX_IMAGE_WIDTH.

        @param pixmap: Source pixmap with its native size.
        @return Tuple (width, height) preserving the aspect ratio.
        """
        if pixmap.width() <= _MAX_IMAGE_WIDTH:
            return pixmap.width(), pixmap.height()
        width = _MAX_IMAGE_WIDTH
        height = round(pixmap.height() * width / pixmap.width())
        return width, height

    def _update_buttons(self) -> None:
        """@brief Enable or disable the navigation buttons for the current slide."""
        self._prev_button.setEnabled(self._index > 0)
        self._next_button.setEnabled(self._index < self._slide_count - 1)

    def _show_previous(self) -> None:
        """@brief Slot: navigate to the previous slide."""
        self._show_slide(self._index - 1)

    def _show_next(self) -> None:
        """@brief Slot: navigate to the next slide."""
        self._show_slide(self._index + 1)

    def cleanup(self) -> None:
        """@brief Release the presentation handle and clear the view."""
        self._presentation = None
        self._slide_count = 0
        self._browser.clear()


def _runs_to_html(paragraph) -> str:
    """@brief Convert a text-frame paragraph to HTML keeping inline styles.

    @param paragraph: pptx text frame paragraph object.
    @return HTML-escaped string with <b>/<i>/<u> tags ("" when empty).
    """
    pieces: list[str] = []
    for run in paragraph.runs:
        text = html.escape(run.text)
        if not text:
            continue
        if run.font.bold:
            text = f"<b>{text}</b>"
        if run.font.italic:
            text = f"<i>{text}</i>"
        if run.font.underline:
            text = f"<u>{text}</u>"
        pieces.append(text)
    return "".join(pieces)


def _table_to_html(table) -> str:
    """@brief Convert a pptx table to a bordered HTML table.

    @param table: pptx.table.Table instance.
    @return HTML fragment with the first _MAX_TABLE_ROWS rows.
    """
    rows_html: list[str] = []
    for row_index, row in enumerate(islice(table.rows, _MAX_TABLE_ROWS)):
        tag = "th" if row_index == 0 else "td"
        cells = "".join(f"<{tag}>{html.escape(cell.text)}</{tag}>" for cell in row.cells)
        rows_html.append(f"<tr>{cells}</tr>")
    return (
        "<table border='1' cellspacing='0' cellpadding='4' "
        "style='border-collapse:collapse; margin:8px auto;'>"
        f"{''.join(rows_html)}</table>"
    )
